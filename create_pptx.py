"""Generate DDM501 Final Project presentation PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Helpers ──────────────────────────────────────────────────────────

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x00, 0xB4, 0x5A)
ACCENT_RED = RGBColor(0xE0, 0x4F, 0x5F)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide_content(slide, items, left, top, width, height,
                              font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_section_header(slide, number, title):
    # Section number circle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.6), Inches(0.5), Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

    add_textbox(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.8),
                title, font_size=32, color=WHITE, bold=True)

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(12), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def add_card(slide, left, top, width, height, title, items,
             accent_color=ACCENT_BLUE):
    # Card background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(1.5)

    # Title bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    add_textbox(slide, left + Inches(0.15), top + Inches(0.05),
                width - Inches(0.3), Inches(0.4),
                title, font_size=14, color=WHITE, bold=True)

    add_bullet_slide_content(
        slide, items,
        left + Inches(0.2), top + Inches(0.55),
        width - Inches(0.4), height - Inches(0.65),
        font_size=13, color=LIGHT_GRAY, spacing=Pt(4))


# ── SLIDE 1: Title ──────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide)

# Accent bar top
bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "Skin Cancer Classification System",
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.7),
            "End-to-End ML System for Dermoscopic Image Classification",
            font_size=24, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

# Divider
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.8), Inches(4.3), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
            "DDM501 — AI in Production: From Models to Systems",
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
            "Final Project Presentation",
            font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Key info boxes
for i, (label, value) in enumerate([
    ("9 Classes", "ISIC Dataset"),
    ("4 Models", "CNN · ResNet · EfficientNet · ViT"),
    ("Full MLOps", "Train · Deploy · Monitor"),
]):
    x = Inches(2.2 + i * 3.2)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.6), Inches(2.8), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(1)
    add_textbox(slide, x + Inches(0.1), Inches(5.65), Inches(2.6), Inches(0.5),
                label, font_size=18, color=ACCENT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(6.1), Inches(2.6), Inches(0.5),
                value, font_size=12, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)


# ── SLIDE 2: Agenda ─────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
            "Agenda", font_size=36, color=WHITE, bold=True)

agenda_items = [
    ("1", "Problem & Solution", "Business context, 9-class skin cancer detection", ACCENT_BLUE),
    ("2", "Technical Deep Dive", "Architecture, data pipeline, models, training", ACCENT_GREEN),
    ("3", "Deployment & Monitoring", "Docker, API, Prometheus, Grafana", ACCENT_ORANGE),
    ("4", "Responsible AI", "Explainability, fairness, ethics", ACCENT_RED),
    ("5", "Live Demo & Q&A", "End-to-end demonstration", ACCENT_BLUE),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    y = Inches(1.5 + i * 1.1)
    # Number
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(1.9), y - Inches(0.05), Inches(9), Inches(0.4),
                title, font_size=22, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.9), y + Inches(0.35), Inches(9), Inches(0.35),
                desc, font_size=14, color=MED_GRAY)


# ── SLIDE 3: Problem & Solution (Rubric: 15%) ───────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 1, "Problem & Solution")

add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.5),
         "THE PROBLEM", [
             "• Skin cancer is the most common cancer globally",
             "• Early detection dramatically improves survival rates",
             "• Dermatologists face high workload & diagnostic variability",
             "• Visual inspection accuracy varies 65-80%",
         ], ACCENT_RED)

add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.5),
         "OUR SOLUTION", [
             "• AI-assisted dermoscopic image classification",
             "• 9 lesion classes (3 malignant, 6 benign)",
             "• Decision support tool — human-in-the-loop",
             "• Real-time API with confidence & risk assessment",
         ], ACCENT_GREEN)

add_card(slide, Inches(0.5), Inches(4.5), Inches(3.8), Inches(2.5),
         "SUCCESS METRICS", [
             "• Missed malignancy rate < 5%",
             "• API latency p95 < 2 seconds",
             "• Macro recall > 0.75",
             "• Macro ROC-AUC > 0.85",
         ], ACCENT_BLUE)

add_card(slide, Inches(4.7), Inches(4.5), Inches(3.8), Inches(2.5),
         "9 SKIN LESION CLASSES", [
             "Malignant: Melanoma, BCC, SCC",
             "Benign: Nevus, Seb. Keratosis,",
             "  Pigmented BK, Actinic Keratosis,",
             "  Dermatofibroma, Vascular Lesion",
         ], ACCENT_ORANGE)

add_card(slide, Inches(8.9), Inches(4.5), Inches(3.8), Inches(2.5),
         "KEY DESIGN DECISIONS", [
             "• Recall-focused (minimize false negatives)",
             "• Multi-model comparison approach",
             "• Full MLOps lifecycle coverage",
             "• Responsible AI first",
         ], RGBColor(0x9B, 0x59, 0xB6))


# ── SLIDE 4: Architecture (Rubric: Technical Deep Dive 40%) ──────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 2, "System Architecture")

# Architecture flow boxes
components = [
    ("ISIC Dataset\n9 Classes\n2,475 images", Inches(0.3), ACCENT_ORANGE),
    ("Data Pipeline\nAugmentor Balance\nTrain/Val/Test Split", Inches(2.9), ACCENT_BLUE),
    ("Model Training\n4 Architectures\nMLflow Tracking", Inches(5.5), ACCENT_GREEN),
    ("Evaluation\nRecall-focused\nModel Selection", Inches(8.1), RGBColor(0x9B, 0x59, 0xB6)),
    ("FastAPI\nServing\nMulti-model", Inches(10.7), ACCENT_RED),
]

for text, x, color in components:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(2.2), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    for j, line in enumerate(text.split("\n")):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE if j == 0 else LIGHT_GRAY
        p.font.bold = (j == 0)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

# Arrows between components
for x in [Inches(2.5), Inches(5.1), Inches(7.7), Inches(10.3)]:
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, x, Inches(2.15), Inches(0.4), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_BLUE
    arrow.line.fill.background()

# Technology stack table
add_card(slide, Inches(0.3), Inches(3.3), Inches(6), Inches(3.8),
         "TECHNOLOGY STACK", [
             "• ML Framework: PyTorch 2.0+ (MPS/CUDA/CPU)",
             "• Model Hub: torchvision + timm (ViT)",
             "• API: FastAPI + Uvicorn (async, auto-docs)",
             "• Experiment Tracking: MLflow",
             "• Containerization: Docker + Docker Compose",
             "• Monitoring: Prometheus + Grafana",
             "• Explainability: pytorch-grad-cam",
             "• Data Balancing: Augmentor (offline oversampling)",
         ], ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(3.3), Inches(6), Inches(3.8),
         "KEY TRADE-OFFS", [
             "• PyTorch > TensorFlow: timm ecosystem, dynamic graphs",
             "• FastAPI > Flask: async, auto OpenAPI, Pydantic",
             "• Weighted loss + sampling: dual imbalance handling",
             "• AdamW + CosineAnnealing: better generalization",
             "• Docker Compose > K8s: right-sized complexity",
             "• Label smoothing 0.1: reduces overconfidence",
             "• Recall > Accuracy: medical false-negative priority",
         ], ACCENT_GREEN)


# ── SLIDE 5: Data Pipeline ──────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 2, "Data Pipeline & Augmentation")

add_card(slide, Inches(0.3), Inches(1.7), Inches(4), Inches(2.8),
         "DATASET", [
             "• ISIC Skin Cancer (Kaggle)",
             "• 2,475 total images",
             "• 9 classes (highly imbalanced)",
             "• Train/Val split: 85%/15% stratified",
             "• Separate Test folder (118 images)",
         ], ACCENT_BLUE)

add_card(slide, Inches(4.6), Inches(1.7), Inches(4), Inches(2.8),
         "CLASS BALANCING (Augmentor)", [
             "• Offline oversampling to 1000/class",
             "• Rotate, flip, zoom, distortion",
             "• Persistent cache (skip if exists)",
             "• 9,000 balanced training images",
             "• + WeightedRandomSampler at runtime",
         ], ACCENT_GREEN)

add_card(slide, Inches(8.9), Inches(1.7), Inches(4), Inches(2.8),
         "ONLINE AUGMENTATION", [
             "• RandomCrop (256 → 224)",
             "• Horizontal + Vertical flip",
             "• Rotation ±30°, ColorJitter",
             "• GaussianBlur, RandomErasing",
             "• ImageNet normalization",
         ], ACCENT_ORANGE)

add_card(slide, Inches(0.3), Inches(4.8), Inches(12.5), Inches(2.3),
         "DATA FLOW", [
             "Raw Images → Load Paths & Labels → Stratified Split (85/15) → Augmentor Balancing (1000/class) → Online Transforms → WeightedRandomSampler → DataLoader (batch=32)",
             "",
             "Three-layer imbalance strategy:  (1) Augmentor offline oversampling  →  (2) Weighted class loss  →  (3) Weighted random sampling",
             "This ensures every batch has balanced class representation even with originally skewed data (e.g., Nevus 357 vs Dermatofibroma 30)",
         ], ACCENT_BLUE)


# ── SLIDE 6: Models ─────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 2, "Model Architectures")

models = [
    ("Custom CNN", ACCENT_BLUE, [
        "• 4 Conv blocks (3→32→64→128→256)",
        "• BatchNorm + Dropout2d per block",
        "• Global Average Pooling",
        "• Classifier: 256→512→128→9",
        "• LR: 1e-3 (train from scratch)",
        "• ~1.2M parameters",
    ]),
    ("ResNet50", ACCENT_GREEN, [
        "• ImageNet V2 pretrained weights",
        "• Freeze all except layer4 + fc",
        "• Custom head: 2048→512→9",
        "• Dropout 0.5 + 0.3",
        "• LR: 1e-4",
        "• ~24M params (2.4M trainable)",
    ]),
    ("EfficientNet-B0", ACCENT_ORANGE, [
        "• ImageNet V1 pretrained weights",
        "• Freeze except features.8 + classifier",
        "• Custom head: 1280→512→9",
        "• Dropout 0.5 + 0.3",
        "• LR: 1e-4",
        "• ~4.3M params (0.8M trainable)",
    ]),
    ("ViT-B/16", ACCENT_RED, [
        "• timm pretrained (ImageNet-21k)",
        "• Freeze except blocks 10-11 + norm",
        "• Custom head: 768→256→9 (GELU)",
        "• LayerNorm + Dropout 0.3",
        "• LR: 5e-5",
        "• ~86M params (14M trainable)",
    ]),
]

for i, (name, color, items) in enumerate(models):
    x = Inches(0.3 + i * 3.2)
    add_card(slide, x, Inches(1.7), Inches(3), Inches(3.5), name, items, color)

add_card(slide, Inches(0.3), Inches(5.5), Inches(12.5), Inches(1.6),
         "TRAINING CONFIGURATION", [
             "Optimizer: AdamW (weight_decay=1e-4)  |  Scheduler: CosineAnnealingWarmRestarts (T₀=10, T_mult=2)  |  Loss: CrossEntropy (weighted + label_smoothing=0.1)",
             "Early Stopping: patience=7  |  Epochs: 30 max  |  Batch: 32  |  Device: CUDA/MPS/CPU auto-detect  |  Experiment Tracking: MLflow",
         ], ACCENT_BLUE)


# ── SLIDE 7: Results ─────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 2, "Training Results & Evaluation")

add_textbox(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
            "Model comparison on validation set (best checkpoint):",
            font_size=18, color=LIGHT_GRAY)

# Results table header
header_y = Inches(2.2)
cols = [("Model", Inches(0.6), Inches(2.5)),
        ("Best Epoch", Inches(3.1), Inches(1.5)),
        ("Val Accuracy", Inches(4.6), Inches(1.8)),
        ("Val F1", Inches(6.4), Inches(1.5)),
        ("Val Recall", Inches(7.9), Inches(1.5)),
        ("Status", Inches(9.4), Inches(2))]

for text, x, w in cols:
    add_textbox(slide, x, header_y, w, Inches(0.4),
                text, font_size=16, color=ACCENT_BLUE, bold=True)

# Divider
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.65), Inches(11), Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Placeholder rows — update with real results
results_data = [
    ("ResNet50", "TBD", "TBD", "TBD", "TBD", "Best performer"),
    ("EfficientNet-B0", "TBD", "TBD", "TBD", "TBD", "Default deploy"),
    ("ViT-B/16", "TBD", "TBD", "TBD", "TBD", "Transformer baseline"),
    ("Custom CNN", "TBD", "TBD", "TBD", "TBD", "Scratch baseline"),
]

for i, (model, epoch, acc, f1, recall, status) in enumerate(results_data):
    y = Inches(2.8 + i * 0.45)
    row_color = WHITE if i % 2 == 0 else LIGHT_GRAY
    for text, x, w in zip([model, epoch, acc, f1, recall, status],
                           [c[1] for c in cols], [c[2] for c in cols]):
        add_textbox(slide, x, y, w, Inches(0.4), text, font_size=14, color=row_color)

add_textbox(slide, Inches(0.6), Inches(4.8), Inches(11), Inches(0.5),
            "📊 Update TBD values with actual results after training completes",
            font_size=14, color=ACCENT_ORANGE, bold=True)

add_card(slide, Inches(0.3), Inches(5.5), Inches(6), Inches(1.6),
         "EVALUATION METRICS", [
             "• Accuracy, Precision, Recall, F1 (macro/weighted)",
             "• ROC-AUC (one-vs-rest), Confusion Matrix",
             "• Per-class classification report",
             "• Cross-model comparison (recall-ranked)",
         ], ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(5.5), Inches(6), Inches(1.6),
         "MODEL SELECTION CRITERIA", [
             "• Primary: Macro Recall (minimize missed cancer)",
             "• Secondary: F1-Score (balance precision/recall)",
             "• Clinical priority: false negatives > false positives",
             "• Best model auto-selected for deployment",
         ], ACCENT_GREEN)


# ── SLIDE 8: Deployment & Monitoring ─────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 3, "Deployment & Monitoring")

add_card(slide, Inches(0.3), Inches(1.7), Inches(4), Inches(2.8),
         "DOCKER DEPLOYMENT", [
             "• Multi-stage Dockerfile (slim build)",
             "• Non-root user (appuser)",
             "• Docker Compose: 4 services",
             "  - API (port 8000)",
             "  - Prometheus (port 9090)",
             "  - Grafana (port 3000)",
             "  - MLflow (port 5000)",
         ], ACCENT_BLUE)

add_card(slide, Inches(4.6), Inches(1.7), Inches(4), Inches(2.8),
         "REST API (FastAPI)", [
             "• POST /predict — classify image",
             "• GET /models — list available models",
             "• GET /health — health check",
             "• GET /classes — class metadata",
             "• GET /metrics — Prometheus metrics",
             "• Auto OpenAPI docs at /docs",
             "• Multi-model selection per request",
         ], ACCENT_GREEN)

add_card(slide, Inches(8.9), Inches(1.7), Inches(4), Inches(2.8),
         "API RESPONSE", [
             '• predicted_class: "melanoma"',
             "• confidence: 0.87",
             "• is_malignant: true",
             '• risk_level: "HIGH"',
             "• probabilities: {all 9 classes}",
             "• model_used: resnet50",
             "• Input validation: type + size",
         ], ACCENT_ORANGE)

add_card(slide, Inches(0.3), Inches(4.8), Inches(6), Inches(2.4),
         "PROMETHEUS METRICS", [
             "• skin_cancer_predictions_total (by class, risk)",
             "• skin_cancer_prediction_latency_seconds",
             "• skin_cancer_prediction_confidence",
             "• skin_cancer_malignant_predictions_total",
             "• skin_cancer_low_confidence_predictions_total",
             "• skin_cancer_class_distribution_total (drift)",
         ], RGBColor(0xE6, 0x52, 0x2C))

add_card(slide, Inches(6.8), Inches(4.8), Inches(6), Inches(2.4),
         "GRAFANA DASHBOARDS", [
             "• Prediction volume & rate over time",
             "• Latency distribution (p50/p95/p99)",
             "• Confidence histogram",
             "• Class distribution (data drift detection)",
             "• Malignant vs benign ratio",
             "• Alert rules for anomalies",
         ], ACCENT_GREEN)


# ── SLIDE 9: Testing & CI/CD ────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 3, "Testing & Quality Assurance")

add_card(slide, Inches(0.3), Inches(1.7), Inches(3), Inches(3),
         "UNIT TESTS", [
             "• Model output shapes",
             "• Freeze/unfreeze behavior",
             "• Forward + backward pass",
             "• Gradient flow verification",
             "• Model factory selection",
         ], ACCENT_BLUE)

add_card(slide, Inches(3.6), Inches(1.7), Inches(3), Inches(3),
         "API TESTS", [
             "• Endpoint response codes",
             "• Input validation errors",
             "• File type/size rejection",
             "• OpenAPI spec availability",
             "• Health check endpoint",
         ], ACCENT_GREEN)

add_card(slide, Inches(6.9), Inches(1.7), Inches(3), Inches(3),
         "DATA QUALITY TESTS", [
             "• Transform output shapes",
             "• Dataset item integrity",
             "• Class constant consistency",
             "• Class weight computation",
             "• Weighted sampler behavior",
         ], ACCENT_ORANGE)

add_card(slide, Inches(10.2), Inches(1.7), Inches(2.8), Inches(3),
         "INTEGRATION TESTS", [
             "• End-to-end prediction",
             "• Result schema validation",
             "• Probability sum check",
             "• Risk level output",
             "• Multi-model inference",
         ], ACCENT_RED)

add_card(slide, Inches(0.3), Inches(5.0), Inches(12.5), Inches(2.2),
         "MLOps PRACTICES", [
             "• MLflow Experiment Tracking: all hyperparameters, metrics, and model artifacts logged per training run",
             "• Model Versioning: best checkpoint saved with epoch, metrics, and architecture metadata",
             "• Reproducibility: fixed random seeds, stratified splits, pinned dependencies (requirements.txt)",
             "• Security: non-root Docker, multi-stage builds, read-only model volumes, input validation, no PII",
         ], RGBColor(0x9B, 0x59, 0xB6))


# ── SLIDE 10: Responsible AI (Rubric: 15%) ───────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 4, "Responsible AI")

add_card(slide, Inches(0.3), Inches(1.7), Inches(4), Inches(3),
         "EXPLAINABILITY", [
             "• Grad-CAM heatmaps for CNN/ResNet/EfficientNet",
             "• ViT-compatible attention visualization",
             "• Visual explanation of model decisions",
             "• Helps clinicians understand predictions",
             "• Batch explainability report generation",
         ], ACCENT_BLUE)

add_card(slide, Inches(4.6), Inches(1.7), Inches(4), Inches(3),
         "FAIRNESS ANALYSIS", [
             "• Per-class recall/precision disparity metrics",
             "• Best/worst class gap analysis",
             "• Equalized odds proxy check (gap < 0.2)",
             "• Optional subgroup auditing",
             "• Recall-focused: minimize missed diagnosis",
         ], ACCENT_GREEN)

add_card(slide, Inches(8.9), Inches(1.7), Inches(4), Inches(3),
         "ETHICAL CONSIDERATIONS", [
             "• Decision SUPPORT — not autonomous diagnosis",
             "• Human-in-the-loop design philosophy",
             "• Risk level flagging for malignant cases",
             "• Confidence threshold for low-certainty alerts",
             "• Generated ethics report artifact",
         ], ACCENT_RED)

add_card(slide, Inches(0.3), Inches(5.0), Inches(6), Inches(2.2),
         "DATA PRIVACY", [
             "• ISIC images are de-identified (no PII)",
             "• No patient data stored or transmitted",
             "• API processes images in-memory only",
             "• HIPAA-aware design considerations",
         ], ACCENT_ORANGE)

add_card(slide, Inches(6.8), Inches(5.0), Inches(6), Inches(2.2),
         "GOVERNANCE & RISK", [
             "• Model limitations clearly documented",
             "• Not intended for clinical deployment without validation",
             "• Continuous monitoring for prediction drift",
             "• Generated responsible AI report with recommendations",
         ], RGBColor(0x9B, 0x59, 0xB6))


# ── SLIDE 11: Live Demo ─────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 5, "Live Demo")

demo_steps = [
    ("1", "Start Services", "docker-compose up → API + Prometheus + Grafana + MLflow", ACCENT_BLUE),
    ("2", "API Prediction", "POST /predict with dermoscopic image → class + confidence + risk", ACCENT_GREEN),
    ("3", "Multi-Model", "Switch models via ?model_name=resnet50 / efficientnet / vit", ACCENT_ORANGE),
    ("4", "Monitoring", "Grafana dashboards: latency, class distribution, alerts", ACCENT_RED),
    ("5", "Explainability", "Grad-CAM heatmap showing model attention regions", RGBColor(0x9B, 0x59, 0xB6)),
]

for i, (num, title, desc, color) in enumerate(demo_steps):
    y = Inches(1.7 + i * 1.05)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.55), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(1.6), y - Inches(0.05), Inches(10), Inches(0.35),
                title, font_size=20, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.6), y + Inches(0.3), Inches(10), Inches(0.35),
                desc, font_size=14, color=LIGHT_GRAY)


# ── SLIDE 12: Thank You ─────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.2),
            "Thank You",
            font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.7),
            "Questions & Discussion",
            font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.2), Inches(3.3), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
            "DDM501 — AI in Production: From Models to Systems",
            font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
            "Skin Cancer Classification — End-to-End ML System",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key links
for i, (label, url) in enumerate([
    ("GitHub Repository", "github.com/..."),
    ("API Docs", "localhost:8000/docs"),
    ("Grafana Dashboard", "localhost:3000"),
]):
    x = Inches(2.5 + i * 3)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6), Inches(2.5), Inches(0.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(1)
    add_textbox(slide, x + Inches(0.1), Inches(6.05), Inches(2.3), Inches(0.35),
                label, font_size=13, color=ACCENT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(6.35), Inches(2.3), Inches(0.35),
                url, font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────

output_path = "DDM501_Final_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
